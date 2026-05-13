import os
import sys
import json
from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.text import PP_ALIGN
from pptx.util import Emu, Inches, Pt

sys.path.insert(0, os.path.dirname(os.path.abspath(__file__)))
from lib_house import RESULTS_DIR, FIGURES_DIR

CSUB_BLUE = RGBColor(0x00, 0x35, 0x94)
CSUB_GOLD = RGBColor(0xFD, 0xB9, 0x13)
WHITE     = RGBColor(0xFF, 0xFF, 0xFF)
INK       = RGBColor(0x20, 0x25, 0x35)
MUTED     = RGBColor(0x60, 0x68, 0x78)
PANEL     = RGBColor(0xF5, 0xF7, 0xFB)
GOLD_TINT = RGBColor(0xFF, 0xF4, 0xD0)

REPO  = os.path.dirname(os.path.dirname(os.path.abspath(__file__)))
FIG   = FIGURES_DIR

SLIDE_W = Inches(13.333)
SLIDE_H = Inches(7.5)


def add_bg(slide, color):
    bg = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, SLIDE_W, SLIDE_H)
    bg.line.fill.background()
    bg.fill.solid()
    bg.fill.fore_color.rgb = color
    sp = bg._element
    sp.getparent().remove(sp)
    slide.shapes._spTree.insert(2, sp)
    return bg


def add_bar(slide, x, y, w, h, color):
    bar = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, x, y, w, h)
    bar.line.fill.background()
    bar.fill.solid()
    bar.fill.fore_color.rgb = color
    return bar


def add_text(slide, x, y, w, h, text, *, size=18, bold=False, color=INK,
             align=PP_ALIGN.LEFT, font="Calibri"):
    tb = slide.shapes.add_textbox(x, y, w, h)
    tf = tb.text_frame
    tf.word_wrap = True
    tf.margin_left = Emu(0); tf.margin_right = Emu(0)
    tf.margin_top = Emu(0); tf.margin_bottom = Emu(0)
    lines = text.split("\n") if isinstance(text, str) else text
    for i, ln in enumerate(lines):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.alignment = align
        r = p.add_run()
        r.text = ln
        r.font.name = font
        r.font.size = Pt(size)
        r.font.bold = bold
        r.font.color.rgb = color
    return tb


def header(slide, title, subtitle=None):
    add_bar(slide, Inches(0), Inches(0.35), Inches(0.25), Inches(0.9), CSUB_GOLD)
    add_text(slide, Inches(0.55), Inches(0.3), Inches(12.2), Inches(0.8),
             title, size=30, bold=True, color=CSUB_BLUE)
    if subtitle:
        add_text(slide, Inches(0.55), Inches(1.05), Inches(12.2), Inches(0.5),
                 subtitle, size=16, color=MUTED)
    add_bar(slide, Inches(0), Inches(7.35), SLIDE_W, Inches(0.15), CSUB_BLUE)


# ----- slides --------------------------------------------------------------

def slide_title(prs):
    s = prs.slides.add_slide(prs.slide_layouts[6])
    add_bg(s, CSUB_BLUE)
    add_text(s, Inches(0.7), Inches(1.7), Inches(12.0), Inches(2.5),
             "House Price Prediction\nfrom Specs and Exterior Photos",
             size=50, bold=True, color=WHITE)
    add_bar(s, Inches(0.7), Inches(4.4), Inches(2.0), Inches(0.08), CSUB_GOLD)
    add_text(s, Inches(0.7), Inches(4.6), Inches(12.0), Inches(0.6),
             "CSU Bakersfield, CMPS 4450 Data Mining Final",
             size=22, color=CSUB_GOLD)
    add_text(s, Inches(0.7), Inches(5.3), Inches(12.0), Inches(0.6),
             "Noah Gallego",
             size=18, color=WHITE)


def slide_question(prs):
    s = prs.slides.add_slide(prs.slide_layouts[6])
    add_bg(s, WHITE)
    header(s, "The question",
           "specs already get us most of the way. can the photo close the gap?")
    add_text(s, Inches(0.55), Inches(2.0), Inches(12.2), Inches(0.5),
             "what we know about each house", size=20, bold=True, color=CSUB_BLUE)
    add_text(s, Inches(0.55), Inches(2.55), Inches(12.2), Inches(1.5),
             "beds, baths, square footage, city, and one exterior photo.",
             size=20, color=INK)
    add_text(s, Inches(0.55), Inches(4.3), Inches(12.2), Inches(0.5),
             "the experiment", size=20, bold=True, color=CSUB_BLUE)
    add_text(s, Inches(0.55), Inches(4.85), Inches(12.2), Inches(2.0),
             "fit a model on specs alone. then add features from the photo.\n"
             "does the second model beat the first?",
             size=20, color=INK)


def slide_data(prs):
    s = prs.slides.add_slide(prs.slide_layouts[6])
    add_bg(s, WHITE)
    header(s, "Data",
           "15,474 southern california listings, $195k to $2M")
    s.shapes.add_picture(os.path.join(FIG, "fig01_price_hist.png"),
                         Inches(0.4), Inches(1.7), width=Inches(6.5))
    s.shapes.add_picture(os.path.join(FIG, "fig03_top_cities.png"),
                         Inches(7.0), Inches(1.7), width=Inches(6.0))


def slide_samples(prs):
    s = prs.slides.add_slide(prs.slide_layouts[6])
    add_bg(s, WHITE)
    header(s, "What a row looks like",
           "one photo plus a few numeric fields per house")
    s.shapes.add_picture(os.path.join(FIG, "fig02_sample_houses.png"),
                         Inches(1.7), Inches(1.55), height=Inches(5.5))


def slide_split(prs):
    s = prs.slides.add_slide(prs.slide_layouts[6])
    add_bg(s, WHITE)
    header(s, "Train, validation, test",
           "how we use the 15,474 listings")

    add_text(s, Inches(0.55), Inches(1.7), Inches(12.2), Inches(0.5),
             "how we set it up", size=20, bold=True, color=CSUB_BLUE)
    add_text(s, Inches(0.55), Inches(2.15), Inches(12.2), Inches(3.5),
             "shuffle the listings and split 65 / 15 / 20 into train, validation, and test\n"
             "fit the models on train\n"
             "pick the model and its settings (tree depth, KNN k) by validation error\n"
             "re-run the whole comparison on 500 random reshuffles to check stability\n"
             "open the test set once, at the end, with every choice already locked",
             size=17, color=INK)

    add_text(s, Inches(0.55), Inches(5.2), Inches(12.2), Inches(0.5),
             "why this design", size=20, bold=True, color=CSUB_BLUE)
    add_text(s, Inches(0.55), Inches(5.65), Inches(12.2), Inches(1.5),
             "one fixed split swings by several thousand dollars just from luck.\n"
             "the reshuffles smooth that out, and the test number stays honest because\n"
             "it never gets to inform a choice.",
             size=16, color=INK)


def slide_why_not_pixels(prs):
    s = prs.slides.add_slide(prs.slide_layouts[6])
    add_bg(s, WHITE)
    header(s, "Why not just feed in the pixels?",
           "the dimensionality wall")

    add_text(s, Inches(0.55), Inches(1.9), Inches(12.2), Inches(0.6),
             "each photo is 128 x 128 x 3 = 49,152 pixel values",
             size=22, bold=True, color=CSUB_BLUE)
    add_text(s, Inches(0.55), Inches(2.7), Inches(12.2), Inches(0.6),
             "training set has roughly 10,000 houses",
             size=22, bold=True, color=CSUB_BLUE)

    add_text(s, Inches(0.55), Inches(4.0), Inches(12.2), Inches(0.6),
             "more inputs than examples", size=20, bold=True, color=INK)
    add_text(s, Inches(0.55), Inches(4.6), Inches(12.2), Inches(2.5),
             "the model can fit anything in training and learn nothing useful.\n"
             "linear regression is underdetermined.\n"
             "trees and KNN see pixel positions as unrelated columns — shift the house\n"
             "ten pixels left and every input value changes.",
             size=17, color=INK)


def slide_filter_intro(prs):
    s = prs.slides.add_slide(prs.slide_layouts[6])
    add_bg(s, WHITE)
    header(s, "A filter is a 3x3 window that slides across the photo",
           "the kernel is fixed, the output is a smaller summary")
    gif_path = os.path.join(REPO, "slides_assets", "sobel_sliding.gif")
    if os.path.exists(gif_path):
        s.shapes.add_picture(gif_path, Inches(1.2), Inches(1.7), height=Inches(4.5))
    add_text(s, Inches(0.55), Inches(6.5), Inches(12.2), Inches(0.7),
             "at each position, multiply the 9 pixels under the window by the kernel weights and sum.\n"
             "that sum becomes one output pixel.  here the kernel detects horizontal edges (sobel-y).",
             size=15, color=INK)


def slide_filter_bank(prs):
    s = prs.slides.add_slide(prs.slide_layouts[6])
    add_bg(s, WHITE)
    header(s, "Our six filters",
           "one photo, six response maps")
    s.shapes.add_picture(os.path.join(FIG, "fig19_filter_bank.png"),
                         Inches(0.2), Inches(1.9), width=Inches(13.0))
    add_text(s, Inches(0.55), Inches(5.5), Inches(12.2), Inches(0.6),
             "edges catch the roofline and window frames.  color dominances pick out sky, lawn, and roof.",
             size=15, color=INK)


def slide_pooling(prs):
    s = prs.slides.add_slide(prs.slide_layouts[6])
    add_bg(s, WHITE)
    header(s, "Pooling: 128 x 128 down to 4 x 4",
           "each filter response collapses to 16 numbers")

    add_text(s, Inches(0.55), Inches(1.9), Inches(12.2), Inches(0.5),
             "the idea", size=20, bold=True, color=CSUB_BLUE)
    add_text(s, Inches(0.55), Inches(2.4), Inches(12.2), Inches(3.0),
             "split each 128 x 128 response into a 4 x 4 grid of 32 x 32 patches.\n"
             "take the mean of each patch.\n"
             "now each filter is 16 numbers instead of 16,384.",
             size=18, color=INK)

    add_text(s, Inches(0.55), Inches(4.8), Inches(12.2), Inches(0.5),
             "what the model sees", size=20, bold=True, color=CSUB_BLUE)
    add_text(s, Inches(0.55), Inches(5.3), Inches(12.2), Inches(2.0),
             "six filters x 16 cells = 96 spatially-aware numbers per photo.\n"
             "plus the 4 specs.  100 features total per house.",
             size=18, color=INK)


def slide_features(prs):
    s = prs.slides.add_slide(prs.slide_layouts[6])
    add_bg(s, WHITE)
    header(s, "What the model sees, end to end",
           "100 numbers per house")

    add_text(s, Inches(0.55), Inches(2.0), Inches(6.0), Inches(0.6),
             "from the listing", size=22, bold=True, color=CSUB_BLUE)
    add_text(s, Inches(0.55), Inches(2.7), Inches(6.0), Inches(4.0),
             "beds\nbaths\nsquare footage\ntypical price in that city",
             size=20, color=INK)

    add_text(s, Inches(7.0), Inches(2.0), Inches(6.0), Inches(0.6),
             "from the photo", size=22, bold=True, color=CSUB_BLUE)
    add_text(s, Inches(7.0), Inches(2.7), Inches(6.0), Inches(4.5),
             "6 filters, each pooled to 4 x 4\n"
             "= 96 numbers describing\n"
             "what's where in the photo",
             size=18, color=INK)


def _styled_table(slide, x, y, w, h, rows, *, header_fill=None, body_fill=None,
                  highlight_rows=(), highlight_fill=None, font_size=13, header_font_size=14):
    if header_fill is None: header_fill = CSUB_BLUE
    if body_fill   is None: body_fill   = WHITE
    if highlight_fill is None: highlight_fill = GOLD_TINT
    n_rows = len(rows); n_cols = len(rows[0])
    tbl_shape = slide.shapes.add_table(n_rows, n_cols, x, y, w, h)
    tbl = tbl_shape.table
    for ci in range(n_cols):
        for ri in range(n_rows):
            cell = tbl.cell(ri, ci)
            cell.text = ""
            tf = cell.text_frame
            p = tf.paragraphs[0]
            p.alignment = PP_ALIGN.LEFT if ci == 0 else PP_ALIGN.CENTER
            r = p.add_run()
            r.text = str(rows[ri][ci])
            r.font.name = "Calibri"
            r.font.size = Pt(header_font_size if ri == 0 else font_size)
            r.font.bold = (ri == 0)
            if ri == 0:
                r.font.color.rgb = WHITE
                cell.fill.solid(); cell.fill.fore_color.rgb = header_fill
            else:
                r.font.color.rgb = INK
                cell.fill.solid()
                if ri in highlight_rows:
                    cell.fill.fore_color.rgb = highlight_fill
                else:
                    cell.fill.fore_color.rgb = body_fill if ri % 2 == 1 else PANEL
    return tbl


def slide_columns(prs):
    s = prs.slides.add_slide(prs.slide_layouts[6])
    add_bg(s, WHITE)
    header(s, "Eight columns from Kaggle",
           "what arrived in the CSV, and what we do with each one")

    rows = [
        ["column",      "example",                "action"],
        ["image_id",    "0",                      "drop  (just a join key)"],
        ["street",      "1317 Van Buren Avenue",  "extract suffix, target-encode"],
        ["citi",        "Salton City, CA",        "target-encode (415 cities)"],
        ["n_citi",      "317",                    "drop  (label-encoded citi, redundant)"],
        ["bed",         "3",                      "keep  (z-score for KNN)"],
        ["bath",        "2.0",                    "keep  (z-score for KNN)"],
        ["sqft",        "1,560",                  "keep  (z-score for KNN)"],
        ["price",       "$201,900",               "target  (what we predict)"],
    ]
    _styled_table(s, Inches(0.55), Inches(1.7), Inches(12.2), Inches(4.8),
                  rows, font_size=14, header_font_size=15)

    add_text(s, Inches(0.55), Inches(6.55), Inches(12.2), Inches(0.6),
             "five engineered features go into every model: bed, bath, sqft, city target mean, street suffix mean.",
             size=14, color=INK)


def slide_encoding(prs):
    s = prs.slides.add_slide(prs.slide_layouts[6])
    add_bg(s, WHITE)
    header(s, "Target encoding for city and street",
           "turn a text column into one number by averaging price within each value")

    # left: city encoding
    add_text(s, Inches(0.55), Inches(1.8), Inches(6.0), Inches(0.5),
             "city  (415 unique values)", size=18, bold=True, color=CSUB_BLUE)
    city_rows = [
        ["city",            "train-mean price"],
        ["Beverly Hills",   "$1,650,000"],
        ["San Diego",       "$  710,000"],
        ["Bakersfield",     "$  295,000"],
        ["…",               "…"],
    ]
    _styled_table(s, Inches(0.55), Inches(2.3), Inches(6.0), Inches(2.2),
                  city_rows, font_size=13, header_font_size=14)
    add_text(s, Inches(0.55), Inches(4.7), Inches(6.0), Inches(2.0),
             "for each row, look up the city,\nuse that average as the feature.\n"
             "cities not seen in train fall back\nto the overall mean price.",
             size=14, color=INK)

    # right: street encoding
    add_text(s, Inches(7.0), Inches(1.8), Inches(6.0), Inches(0.5),
             "street  (extract suffix first)", size=18, bold=True, color=CSUB_BLUE)
    street_rows = [
        ["suffix",   "train-mean price"],
        ["circle",   "$791,000"],
        ["lane",     "$716,000"],
        ["road",     "$746,000"],
        ["boulevard","$629,000"],
    ]
    _styled_table(s, Inches(7.0), Inches(2.3), Inches(6.0), Inches(2.2),
                  street_rows, font_size=13, header_font_size=14)
    add_text(s, Inches(7.0), Inches(4.7), Inches(6.0), Inches(2.0),
             "the 12,401 unique addresses have\nno signal on their own.  the last word\n"
             "(circle, lane, blvd, ...) does.\n11 suffix types, target-encoded.",
             size=14, color=INK)


def slide_normalization(prs):
    s = prs.slides.add_slide(prs.slide_layouts[6])
    add_bg(s, WHITE)
    header(s, "Normalizing the numeric columns",
           "without this, sqft and city mean dominate every distance KNN computes")

    rows = [
        ["feature",            "raw range",          "after z-score"],
        ["bed",                "1 to 6",             "≈ -2 to +3"],
        ["bath",               "1 to 5",             "≈ -2 to +3"],
        ["sqft",               "500 to 5,400",       "≈ -2 to +5"],
        ["city target mean",   "$200k to $1.65M",    "≈ -2 to +4"],
        ["street suffix mean", "$629k to $792k",     "≈ -2 to +3"],
    ]
    _styled_table(s, Inches(0.55), Inches(1.9), Inches(12.2), Inches(3.6),
                  rows, font_size=14, header_font_size=15)

    add_text(s, Inches(0.55), Inches(5.8), Inches(12.2), Inches(0.5),
             "z-score:  (value - train_mean) / train_std",
             size=18, bold=True, color=CSUB_BLUE)
    add_text(s, Inches(0.55), Inches(6.35), Inches(12.2), Inches(0.8),
             "linear and tree don't care about scale. KNN does — its distance is the sum of squared differences.\n"
             "without normalizing, every neighbor lookup is decided by sqft alone.",
             size=14, color=INK)


def slide_three_models(prs):
    s = prs.slides.add_slide(prs.slide_layouts[6])
    add_bg(s, WHITE)
    header(s, "Three models",
           "fit each on the five engineered features")

    ys = Inches(2.0); col_w = Inches(4.0); gap = Inches(0.3)
    xs = [Inches(0.55), Inches(0.55) + col_w + gap, Inches(0.55) + 2 * (col_w + gap)]
    titles = ["linear regression", "decision tree", "k-nearest neighbors"]
    bodies = [
        "one weight per feature.\nfast, transparent.\nstruggles with non-linear\nrelationships like the\nsqft / price curve.",
        "splits the data on one\nfeature at a time.\nhandles non-linear cuts,\ndoesn't care about scale,\ntuned by depth.",
        "predict by averaging the\nk most similar houses.\nneeds normalized\nfeatures, tuned by k and\ndistance function.",
    ]
    for x, t, b in zip(xs, titles, bodies):
        add_bar(s, x, ys, col_w, Inches(0.55), CSUB_BLUE)
        add_text(s, x + Inches(0.15), ys + Inches(0.07), col_w, Inches(0.5),
                 t, size=18, bold=True, color=WHITE)
        add_text(s, x + Inches(0.15), ys + Inches(0.85), col_w - Inches(0.3), Inches(3.5),
                 b, size=15, color=INK)


def slide_knn_distances(prs):
    s = prs.slides.add_slide(prs.slide_layouts[6])
    add_bg(s, WHITE)
    header(s, "KNN, four distance functions",
           "each box summarizes 200 random reshuffles at that k")
    img = os.path.join(FIG, "fig20_knn_distances.png")
    if os.path.exists(img):
        s.shapes.add_picture(img, Inches(0.4), Inches(1.6), width=Inches(9.5))
    add_text(s, Inches(10.1), Inches(1.9), Inches(3.0), Inches(0.5),
             "reading the plot", size=18, bold=True, color=CSUB_BLUE)
    add_text(s, Inches(10.1), Inches(2.4), Inches(3.0), Inches(4.5),
             "each panel is a different\ndistance function.\n\n"
             "each box: spread of\nvalidation MAE at that k.\n\n"
             "gold line connects the\nmedians.\n\n"
             "lower is better.",
             size=13, color=INK)


def slide_baseline_box(prs):
    s = prs.slides.add_slide(prs.slide_layouts[6])
    add_bg(s, WHITE)
    header(s, "Specs only is already strong",
           "error in dollars, lower is better")
    img = os.path.join(FIG, "fig05_baseline_boxplot.png")
    if os.path.exists(img):
        s.shapes.add_picture(img, Inches(0.4), Inches(1.7), width=Inches(8.8))
    add_text(s, Inches(9.5), Inches(2.0), Inches(3.5), Inches(0.6),
             "decision tree wins", size=20, bold=True, color=CSUB_BLUE)
    add_text(s, Inches(9.5), Inches(2.7), Inches(3.5), Inches(4.0),
             "around $145k off,\non average,\nacross many reshuffles.\n\n"
             "linear lags. KNN sits\nbetween the two.",
             size=16, color=INK)


def slide_worst_misses(prs):
    s = prs.slides.add_slide(prs.slide_layouts[6])
    add_bg(s, WHITE)
    header(s, "Where the baseline misses",
           "the top six errors split into two kinds")
    s.shapes.add_picture(os.path.join(FIG, "fig07_worst_misses.png"),
                         Inches(1.8), Inches(1.55), height=Inches(5.6))


def slide_cap_problem(prs):
    s = prs.slides.add_slide(prs.slide_layouts[6])
    add_bg(s, WHITE)
    header(s, "The cap problem",
           "the dataset has a ceiling")
    s.shapes.add_picture(os.path.join(FIG, "fig18_cap_cluster.png"),
                         Inches(1.4), Inches(1.7), width=Inches(10.5))
    add_text(s, Inches(0.55), Inches(6.7), Inches(12.2), Inches(0.5),
             "38 houses share the exact price $1,995,000. the model can't reach the cap, so those listings stay at the top of the error list.",
             size=14, color=INK)


def slide_hypothesis(prs):
    s = prs.slides.add_slide(prs.slide_layouts[6])
    add_bg(s, WHITE)
    header(s, "Three things the photo could fix",
           "the case for adding image features")

    ys = Inches(2.2); col_w = Inches(4.0); gap = Inches(0.3)
    xs = [Inches(0.55), Inches(0.55) + col_w + gap, Inches(0.55) + 2 * (col_w + gap)]
    titles = ["condition", "size beyond sqft", "neighborhood look"]
    bodies = [
        "same specs, different\nupkeep and landscaping.\nthat lives in the photo.",
        "story count, lot size,\nyard space. specs miss\nthose. the photo doesn't.",
        "trees, driveway, fencing,\nlighting. signals the\nstreet, not just the city.",
    ]
    for x, t, b in zip(xs, titles, bodies):
        add_bar(s, x, ys, col_w, Inches(0.55), CSUB_BLUE)
        add_text(s, x + Inches(0.15), ys + Inches(0.07), col_w, Inches(0.5),
                 t, size=20, bold=True, color=WHITE)
        add_text(s, x + Inches(0.15), ys + Inches(0.85), col_w - Inches(0.3), Inches(3.0),
                 b, size=17, color=INK)


def slide_image_box(prs):
    s = prs.slides.add_slide(prs.slide_layouts[6])
    add_bg(s, WHITE)
    header(s, "Adding photo features",
           "same three models, with and without the photo")
    img = os.path.join(FIG, "fig11_specs_vs_image_boxplot.png")
    if os.path.exists(img):
        s.shapes.add_picture(img, Inches(0.4), Inches(1.7), height=Inches(4.6))
    add_text(s, Inches(0.55), Inches(6.5), Inches(12.2), Inches(0.6),
             "no model improved. KNN got much worse.",
             size=22, bold=True, color=CSUB_BLUE)


def slide_bad_photos(prs):
    s = prs.slides.add_slide(prs.slide_layouts[6])
    add_bg(s, WHITE)
    header(s, "Why? Some photos aren't of the house",
           "examples pulled from the worst-miss list")
    s.shapes.add_picture(os.path.join(FIG, "fig16_bad_photos.png"),
                         Inches(0.3), Inches(2.2), width=Inches(12.7))
    add_text(s, Inches(0.55), Inches(5.6), Inches(12.2), Inches(0.6),
             "if the photo shows a harbor or a tree, average color and edge density describe those, not the house.",
             size=16, color=INK)


def slide_feature_scales(prs):
    s = prs.slides.add_slide(prs.slide_layouts[6])
    add_bg(s, WHITE)
    header(s, "Why? Feature scales are wildly different",
           "raw values span orders of magnitude")
    s.shapes.add_picture(os.path.join(FIG, "fig17_feature_scales.png"),
                         Inches(0.3), Inches(1.55), height=Inches(5.0))
    add_text(s, Inches(0.55), Inches(6.7), Inches(12.2), Inches(0.5),
             "pixel variance dwarfs everything else. in a distance-based model that one feature drowns the rest.",
             size=15, color=INK)


def slide_preprocess(prs):
    s = prs.slides.add_slide(prs.slide_layouts[6])
    add_bg(s, WHITE)
    header(s, "What we tried",
           "three preprocessing fixes, each motivated by a problem above")

    rows = [
        ("center crop",          "throw out the outer 20% on every side, focus on what's in the middle of the frame.",
                                 "addresses photos where sky or pavement dominates."),
        ("per-image normalize",  "rescale each photo so its mean brightness matches every other photo.",
                                 "addresses photos shot at different times of day or different exposures."),
        ("add HSV channels",     "compute hue, saturation, and value averages on top of RGB.",
                                 "color signal that doesn't change with brightness."),
    ]
    y = Inches(2.0)
    for name, what, why in rows:
        add_text(s, Inches(0.55), y, Inches(3.0), Inches(0.5),
                 name, size=20, bold=True, color=CSUB_BLUE)
        add_text(s, Inches(3.8), y, Inches(9.0), Inches(0.5),
                 what, size=15, color=INK)
        add_text(s, Inches(3.8), y + Inches(0.45), Inches(9.0), Inches(0.5),
                 why, size=14, color=MUTED)
        y = y + Inches(1.3)


def slide_variant_table(prs, variants):
    s = prs.slides.add_slide(prs.slide_layouts[6])
    add_bg(s, WHITE)
    header(s, "Did preprocessing help?",
           "validation MAE, decision tree, by feature variant")

    left = Inches(1.3); top = Inches(2.0)
    widths = [Inches(4.6), Inches(2.4), Inches(2.4), Inches(2.4)]
    cell_h = Inches(0.55)
    heads = ["feature variant", "MAE", "within $50k", "within $100k"]

    def cell(x, y, w, h, text, *, fill, color, bold=False, size=14):
        shp = s.shapes.add_shape(MSO_SHAPE.RECTANGLE, x, y, w, h)
        shp.line.color.rgb = RGBColor(0xE0, 0xE0, 0xE0)
        shp.fill.solid(); shp.fill.fore_color.rgb = fill
        tf = shp.text_frame
        tf.margin_left = Inches(0.12); tf.margin_right = Inches(0.12)
        tf.margin_top = Inches(0.04); tf.margin_bottom = Inches(0.04)
        p = tf.paragraphs[0]; p.alignment = PP_ALIGN.CENTER
        r = p.add_run(); r.text = text
        r.font.name = "Calibri"; r.font.size = Pt(size); r.font.bold = bold
        r.font.color.rgb = color

    x = left
    for w, h in zip(widths, heads):
        cell(x, top, w, cell_h, h, fill=CSUB_BLUE, color=WHITE, bold=True, size=15)
        x += w

    keys = [
        ("specs only",                "specs_only_tree"),
        ("specs + raw photo",         "specs_+_raw_image_tree"),
        ("specs + center crop",       "specs_+_crop_tree"),
        ("specs + brightness norm",   "specs_+_per-img_norm_tree"),
        ("specs + crop + norm + HSV", "specs_+_crop+norm+HSV_tree"),
    ]
    best_mae = min(variants[k]["mae"] for _, k in keys)
    y = top + cell_h
    for i, (label, k) in enumerate(keys):
        m = variants[k]
        bg = WHITE if i % 2 == 0 else PANEL
        if m["mae"] == best_mae:
            bg = GOLD_TINT
        x = left
        cell(x, y, widths[0], cell_h, label,                  fill=bg, color=INK, size=14); x += widths[0]
        cell(x, y, widths[1], cell_h, f"${m['mae']/1000:.0f}k", fill=bg, color=INK, size=14); x += widths[1]
        cell(x, y, widths[2], cell_h, f"{m['within_50k']:.1%}", fill=bg, color=INK, size=14); x += widths[2]
        cell(x, y, widths[3], cell_h, f"{m['within_100k']:.1%}",fill=bg, color=INK, size=14)
        y += cell_h

    add_text(s, Inches(0.55), Inches(5.9), Inches(12.2), Inches(0.6),
             "no variant beats specs only.",
             size=22, bold=True, color=CSUB_BLUE)
    add_text(s, Inches(0.55), Inches(6.55), Inches(12.2), Inches(0.6),
             "color and edge statistics from a single low-resolution photo can't tell us condition or value.",
             size=15, color=INK)


def slide_spatial_results(prs, spatial):
    s = prs.slides.add_slide(prs.slide_layouts[6])
    add_bg(s, WHITE)
    header(s, "Did spatial features help?",
           "validation MAE, by feature set and model")

    left = Inches(0.6); top = Inches(2.0)
    widths = [Inches(4.8), Inches(2.2), Inches(2.2), Inches(2.0), Inches(2.0)]
    cell_h = Inches(0.55)
    heads = ["feature set", "linear", "tree", "knn", "best"]

    def cell(x, y, w, h, text, *, fill, color, bold=False, size=14):
        shp = s.shapes.add_shape(MSO_SHAPE.RECTANGLE, x, y, w, h)
        shp.line.color.rgb = RGBColor(0xE0, 0xE0, 0xE0)
        shp.fill.solid(); shp.fill.fore_color.rgb = fill
        tf = shp.text_frame
        tf.margin_left = Inches(0.12); tf.margin_right = Inches(0.12)
        tf.margin_top = Inches(0.04); tf.margin_bottom = Inches(0.04)
        p = tf.paragraphs[0]; p.alignment = PP_ALIGN.CENTER
        r = p.add_run(); r.text = text
        r.font.name = "Calibri"; r.font.size = Pt(size); r.font.bold = bold
        r.font.color.rgb = color

    x = left
    for w, h in zip(widths, heads):
        cell(x, top, w, cell_h, h, fill=CSUB_BLUE, color=WHITE, bold=True, size=14)
        x += w

    rows = [
        ("specs only",            "specs_only"),
        ("specs + spatial (96)",  "specs_+_spatial_96"),
        ("spatial only",          "spatial_only_no_specs"),
    ]
    y = top + cell_h
    for i, (label, prefix) in enumerate(rows):
        lin  = spatial[f"{prefix}_linear"]["mae"]
        tree = spatial[f"{prefix}_tree"]["mae"]
        knn  = spatial[f"{prefix}_knn"]["mae"]
        best = min(lin, tree, knn)
        bg = WHITE if i % 2 == 0 else PANEL
        if best == min(spatial[f"specs_only_tree"]["mae"], spatial[f"specs_+_spatial_96_tree"]["mae"], spatial[f"spatial_only_no_specs_tree"]["mae"]) and prefix == "specs_only":
            bg = GOLD_TINT
        x = left
        cell(x, y, widths[0], cell_h, label, fill=bg, color=INK, size=14); x += widths[0]
        cell(x, y, widths[1], cell_h, f"${lin/1000:.0f}k",  fill=bg, color=INK, size=14); x += widths[1]
        cell(x, y, widths[2], cell_h, f"${tree/1000:.0f}k", fill=bg, color=INK, size=14); x += widths[2]
        cell(x, y, widths[3], cell_h, f"${knn/1000:.0f}k",  fill=bg, color=INK, size=14); x += widths[3]
        cell(x, y, widths[4], cell_h, f"${best/1000:.0f}k", fill=bg, color=INK, bold=True, size=14)
        y += cell_h

    add_text(s, Inches(0.55), Inches(5.7), Inches(12.2), Inches(0.6),
             "spatial features alone get to $284k — barely better than guessing the average.",
             size=16, color=INK)
    add_text(s, Inches(0.55), Inches(6.25), Inches(12.2), Inches(0.6),
             "adding them on top of specs makes tree and KNN worse, not better.",
             size=16, color=INK)
    add_text(s, Inches(0.55), Inches(6.85), Inches(12.2), Inches(0.5),
             "specs still win.",
             size=18, bold=True, color=CSUB_BLUE)


def slide_subgroup(prs):
    s = prs.slides.add_slide(prs.slide_layouts[6])
    add_bg(s, WHITE)
    header(s, "Per-city, with and without the photo",
           "two cities benefit a little. the rest get worse.")
    s.shapes.add_picture(os.path.join(FIG, "fig12_per_city_delta.png"),
                         Inches(0.6), Inches(1.6), height=Inches(5.5))


def slide_kmeans(prs):
    s = prs.slides.add_slide(prs.slide_layouts[6])
    add_bg(s, WHITE)
    header(s, "K-means on photo features alone",
           "if photos carried price signal, clusters would land at different prices")
    s.shapes.add_picture(os.path.join(FIG, "fig14_kmeans_image.png"),
                         Inches(0.3), Inches(1.6), width=Inches(9.0))
    add_text(s, Inches(9.6), Inches(2.0), Inches(3.4), Inches(0.6),
             "what we found", size=20, bold=True, color=CSUB_BLUE)
    add_text(s, Inches(9.6), Inches(2.7), Inches(3.4), Inches(4.0),
             "four clusters, mean\nprices land within\n$80k of each other.\n\n"
             "the photo features\ndon't separate price.",
             size=16, color=INK)


def slide_test_table(prs, m):
    s = prs.slides.add_slide(prs.slide_layouts[6])
    add_bg(s, WHITE)
    header(s, "Final test results",
           "trained on train + val, evaluated once on the 3,095 held-out test rows")

    left = Inches(0.7); top = Inches(1.9)
    widths = [Inches(3.4), Inches(2.4), Inches(1.7), Inches(1.7), Inches(1.5), Inches(1.5)]
    cell_h = Inches(0.5)
    heads = ["model", "features", "MAE", "RMSE", "<$50k", "<$100k"]

    def cell(x, y, w, h, text, *, fill, color, bold=False, size=13):
        shp = s.shapes.add_shape(MSO_SHAPE.RECTANGLE, x, y, w, h)
        shp.line.color.rgb = RGBColor(0xE0, 0xE0, 0xE0)
        shp.fill.solid(); shp.fill.fore_color.rgb = fill
        tf = shp.text_frame
        tf.margin_left = Inches(0.1); tf.margin_right = Inches(0.1)
        tf.margin_top = Inches(0.03); tf.margin_bottom = Inches(0.03)
        p = tf.paragraphs[0]; p.alignment = PP_ALIGN.CENTER
        r = p.add_run(); r.text = text
        r.font.name = "Calibri"; r.font.size = Pt(size); r.font.bold = bold
        r.font.color.rgb = color

    x = left
    for w, h in zip(widths, heads):
        cell(x, top, w, cell_h, h, fill=CSUB_BLUE, color=WHITE, bold=True, size=14)
        x += w

    rows = [
        ("constant",                  "none",            m["test_constant_predictor"]),
        ("linear regression",         "specs",           m["test_linear_specs"]),
        ("linear regression",         "specs + spatial", m["test_linear_specs_spatial"]),
        ("decision tree, d=15",       "specs",           m["test_tree_specs_depth15"]),
        ("decision tree, d=10",       "specs + spatial", m["test_tree_specs_spatial_d10"]),
        ("KNN, k=15",                 "specs",           m["test_knn_specs_k15"]),
        ("KNN, k=15",                 "specs + spatial", m["test_knn_specs_spatial_k15"]),
    ]
    best_idx = min(range(len(rows)), key=lambda i: rows[i][2]["mae"])
    y = top + cell_h
    for i, (name, feat, met) in enumerate(rows):
        bg = WHITE if i % 2 == 0 else PANEL
        if i == best_idx:
            bg = GOLD_TINT
        x = left
        cell(x, y, widths[0], cell_h, name, fill=bg, color=INK, size=12); x += widths[0]
        cell(x, y, widths[1], cell_h, feat, fill=bg, color=INK, size=12); x += widths[1]
        cell(x, y, widths[2], cell_h, f"${met['mae']/1000:.0f}k",  fill=bg, color=INK, size=12); x += widths[2]
        cell(x, y, widths[3], cell_h, f"${met['rmse']/1000:.0f}k", fill=bg, color=INK, size=12); x += widths[3]
        cell(x, y, widths[4], cell_h, f"{met['within_50k']:.1%}",  fill=bg, color=INK, size=12); x += widths[4]
        cell(x, y, widths[5], cell_h, f"{met['within_100k']:.1%}", fill=bg, color=INK, size=12)
        y += cell_h

    add_text(s, Inches(0.55), Inches(6.2), Inches(12.2), Inches(0.6),
             "specs-only decision tree wins on test.",
             size=22, bold=True, color=CSUB_BLUE)
    add_text(s, Inches(0.55), Inches(6.85), Inches(12.2), Inches(0.5),
             "MAE $142k. 33% of guesses within $50k, 54% within $100k.",
             size=15, color=INK)


def slide_conclusion(prs):
    s = prs.slides.add_slide(prs.slide_layouts[6])
    add_bg(s, CSUB_BLUE)
    add_bar(s, Inches(0), Inches(1.2), SLIDE_W, Inches(0.08), CSUB_GOLD)
    add_text(s, Inches(0.6), Inches(0.5), Inches(12.2), Inches(0.7),
             "Conclusion", size=34, bold=True, color=WHITE)

    add_text(s, Inches(0.6), Inches(1.6), Inches(12.2), Inches(1.4),
             "specs already carry the price signal.\n"
             "the photo, hand-summarized in numpy, does not add anything on top.",
             size=24, color=CSUB_GOLD)

    add_text(s, Inches(0.6), Inches(4.0), Inches(12.2), Inches(0.6),
             "what we learned", size=20, bold=True, color=CSUB_GOLD)
    add_text(s, Inches(0.6), Inches(4.55), Inches(12.2), Inches(3.0),
             "feeding raw pixels doesn't work — too many inputs, too few examples,\n"
             "and classical models can't see that adjacent pixels are related.\n\n"
             "we built the principled fix: six hand-designed filters plus 4 x 4 pooling.\n"
             "the model gets spatial info — sky on top, lawn on bottom, edges where the roof is.\n"
             "this still doesn't beat specs.\n\n"
             "the signal exists inside the photo, but recovering it needs filters that LEARN\n"
             "what to look for. that's what CNNs do.  hand-designed kernels can only go so far.",
             size=15, color=WHITE)


def main():
    with open(os.path.join(RESULTS_DIR, "metrics.json")) as f:
        metrics = json.load(f)
    with open(os.path.join(RESULTS_DIR, "spatial_comparison.json")) as f:
        spatial = json.load(f)

    prs = Presentation()
    prs.slide_width  = SLIDE_W
    prs.slide_height = SLIDE_H

    slide_title(prs)
    slide_question(prs)
    slide_data(prs)
    slide_samples(prs)
    slide_split(prs)
    slide_columns(prs)
    slide_encoding(prs)
    slide_normalization(prs)
    slide_three_models(prs)
    slide_knn_distances(prs)
    slide_worst_misses(prs)
    slide_hypothesis(prs)
    slide_why_not_pixels(prs)
    slide_filter_intro(prs)
    slide_filter_bank(prs)
    slide_pooling(prs)
    slide_features(prs)
    slide_spatial_results(prs, spatial)
    slide_subgroup(prs)
    slide_kmeans(prs)
    slide_test_table(prs, metrics)
    slide_conclusion(prs)

    out = os.path.join(REPO, "HousePrice_Final_Project.pptx")
    prs.save(out)
    print(f"saved {out}  ({len(prs.slides)} slides)")


if __name__ == "__main__":
    main()
